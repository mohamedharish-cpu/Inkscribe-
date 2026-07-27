import io
import zipfile
import pdfplumber
from pptx import Presentation

def extract_text_from_pdf(file_source) -> str:
    """Extracts clean text from PDF files."""
    extracted_text = ""
    try:
        with pdfplumber.open(file_source) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    extracted_text += f"\n--- Page {page_num} ---\n{text}"
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return extracted_text

def extract_text_from_pptx(file_source) -> str:
    """Extracts text from PPTX presentations."""
    extracted_text = ""
    try:
        prs = Presentation(file_source)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                extracted_text += f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_text)
    except Exception as e:
        print(f"Error reading PPTX: {e}")
    return extracted_text

def parse_all_uploaded_files(uploaded_files_list) -> dict:
    """
    Parses PDF, PPTX, and ZIP archives seamlessly into text context.
    """
    parsed_documents = {}
    
    for uploaded_file in uploaded_files_list:
        file_name = uploaded_file.name.lower()
        
        # Handle ZIP Archives
        if file_name.endswith('.zip'):
            try:
                with zipfile.ZipFile(uploaded_file) as z:
                    for zip_info in z.infolist():
                        inner_filename = zip_info.filename.lower()
                        if inner_filename.endswith('.pdf'):
                            with z.open(zip_info) as f:
                                pdf_bytes = io.BytesIO(f.read())
                                text = extract_text_from_pdf(pdf_bytes)
                                if text.strip():
                                    parsed_documents[zip_info.filename] = text
                        elif inner_filename.endswith(('.pptx', '.ppt')):
                            with z.open(zip_info) as f:
                                pptx_bytes = io.BytesIO(f.read())
                                text = extract_text_from_pptx(pptx_bytes)
                                if text.strip():
                                    parsed_documents[zip_info.filename] = text
            except Exception as e:
                print(f"Error processing ZIP file {uploaded_file.name}: {e}")
                
        # Handle Direct PDFs
        elif file_name.endswith('.pdf'):
            text = extract_text_from_pdf(uploaded_file)
            if text.strip():
                parsed_documents[uploaded_file.name] = text
                
        # Handle Direct PPTX
        elif file_name.endswith(('.pptx', '.ppt')):
            text = extract_text_from_pptx(uploaded_file)
            if text.strip():
                parsed_documents[uploaded_file.name] = text
            
    return parsed_documents
